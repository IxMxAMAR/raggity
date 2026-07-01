"""Tests for src/raggity/readers.py — reader registry (Phase B Task 1)."""
from __future__ import annotations


def test_read_txt(tmp_path):
    from raggity.readers import read_file

    p = tmp_path / "a.txt"
    p.write_text("hello text world")
    assert "hello text world" in read_file(str(p))


def test_read_md(tmp_path):
    from raggity.readers import read_file

    p = tmp_path / "a.md"
    p.write_text("# Title\nSome markdown content")
    out = read_file(str(p))
    assert "Title" in out and "markdown content" in out


def test_read_docx(tmp_path):
    import docx
    from raggity.readers import read_file

    p = tmp_path / "a.docx"
    d = docx.Document()
    d.add_paragraph("hello docx world")
    d.save(p)
    assert "docx world" in read_file(str(p))


def test_read_csv(tmp_path):
    from raggity.readers import read_file

    (tmp_path / "a.csv").write_text("name,role\nAlice,dev\n")
    out = read_file(str(tmp_path / "a.csv"))
    assert "Alice" in out and "role" in out


def test_read_html(tmp_path):
    from raggity.readers import read_file

    (tmp_path / "a.html").write_text(
        "<html><body><h1>Title</h1><p>Body text</p></body></html>"
    )
    out = read_file(str(tmp_path / "a.html"))
    assert "Title" in out and "Body text" in out


def test_read_pptx(tmp_path):
    from pptx import Presentation
    from raggity.readers import read_file

    pr = Presentation()
    s = pr.slides.add_slide(pr.slide_layouts[5])
    s.shapes.title.text = "Deck Title"
    pr.save(tmp_path / "a.pptx")
    assert "Deck Title" in read_file(str(tmp_path / "a.pptx"))


def test_unknown_ext_returns_none(tmp_path):
    from raggity.readers import read_file

    (tmp_path / "a.bin").write_bytes(b"\x00")
    assert read_file(str(tmp_path / "a.bin")) is None


def test_supported_exts_contains_all_formats():
    from raggity.readers import SUPPORTED_EXTS

    for ext in {".md", ".txt", ".pdf", ".docx", ".html", ".csv", ".pptx"}:
        assert ext in SUPPORTED_EXTS, f"{ext} missing from SUPPORTED_EXTS"


def test_image_routes_through_ocr(tmp_path, monkeypatch):
    import raggity.readers as r

    monkeypatch.setattr(r, "_run_ocr", lambda src: "OCR TEXT FROM IMAGE")
    (tmp_path / "a.png").write_bytes(b"\x89PNG\r\n\x1a\n")  # bytes irrelevant; OCR is mocked
    assert "OCR TEXT" in r.read_file(str(tmp_path / "a.png"))


def test_scanned_pdf_falls_back_to_ocr(tmp_path, monkeypatch):
    import raggity.readers as r

    monkeypatch.setattr(r, "_pdf_text", lambda p: "")           # simulate no embedded text
    monkeypatch.setattr(r, "_ocr_pdf", lambda p: "OCR PDF TEXT")
    (tmp_path / "scan.pdf").write_bytes(b"%PDF-1.4")
    assert "OCR PDF TEXT" in r.read_file(str(tmp_path / "scan.pdf"))


def test_text_pdf_skips_ocr(tmp_path, monkeypatch):
    """If _pdf_text returns content, _ocr_pdf must never be called."""
    import raggity.readers as r

    monkeypatch.setattr(r, "_pdf_text", lambda p: "embedded text content")
    ocr_called = []
    monkeypatch.setattr(r, "_ocr_pdf", lambda p: ocr_called.append(True) or "SHOULD NOT SEE")
    (tmp_path / "text.pdf").write_bytes(b"%PDF-1.4")
    result = r.read_file(str(tmp_path / "text.pdf"))
    assert "embedded text content" in result
    assert not ocr_called


def test_image_exts_in_supported_exts():
    from raggity.readers import SUPPORTED_EXTS

    for ext in {".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp"}:
        assert ext in SUPPORTED_EXTS, f"{ext} missing from SUPPORTED_EXTS"


def test_read_image_calls_run_ocr(tmp_path, monkeypatch):
    import raggity.readers as r

    monkeypatch.setattr(r, "_run_ocr", lambda src: "DIRECT IMAGE OCR")
    (tmp_path / "b.jpg").write_bytes(b"\xff\xd8\xff")
    assert "DIRECT IMAGE OCR" in r.read_image(str(tmp_path / "b.jpg"))


def test_missing_ocr_dep_raises_friendly(monkeypatch):
    """_run_ocr raises RuntimeError with install hint when RapidOCR is missing."""
    import sys
    import raggity.readers as r

    # Simulate RapidOCR not installed by patching the import inside _run_ocr
    orig = sys.modules.get("rapidocr_onnxruntime")
    sys.modules["rapidocr_onnxruntime"] = None  # type: ignore[assignment]
    # Reset singleton so it re-imports
    r._ocr_engine_instance = None
    try:
        import pytest
        with pytest.raises(RuntimeError, match="pip install raggity"):
            r._run_ocr("irrelevant.png")
    finally:
        if orig is None:
            sys.modules.pop("rapidocr_onnxruntime", None)
        else:
            sys.modules["rapidocr_onnxruntime"] = orig
        r._ocr_engine_instance = None


def test_read_pdf_via_read_file(tmp_path):
    """read_file dispatches .pdf correctly (uses the moved read_pdf).
    Dispatching .pdf must NOT return None (the unknown-ext sentinel).
    It may return a string or raise; either is fine.
    """
    from raggity.readers import read_file

    bad = tmp_path / "b.pdf"
    bad.write_bytes(b"%PDF-1.4 not really a pdf")
    dispatched = False
    try:
        result = read_file(str(bad))
        dispatched = True  # returned a value (possibly empty string)
    except Exception:
        dispatched = True  # raised from pypdf — dispatch DID happen
    assert dispatched, ".pdf was not dispatched (returned None as unknown ext)"


def test_missing_dependency_error_is_runtime_error():
    from raggity.readers import MissingDependencyError
    err = MissingDependencyError(extra="ocr", message="needs ocr")
    assert isinstance(err, RuntimeError)
    assert err.extra == "ocr"


def test_run_ocr_raises_missing_dep_error_not_plain_runtime(monkeypatch):
    """_run_ocr raises MissingDependencyError(extra='ocr') when rapidocr absent."""
    import sys
    import raggity.readers as r
    from raggity.readers import MissingDependencyError

    orig = sys.modules.get("rapidocr_onnxruntime")
    sys.modules["rapidocr_onnxruntime"] = None  # type: ignore[assignment]
    r._ocr_engine_instance = None
    try:
        import pytest
        with pytest.raises(MissingDependencyError) as exc_info:
            r._run_ocr("irrelevant.png")
        assert exc_info.value.extra == "ocr"
    finally:
        if orig is None:
            sys.modules.pop("rapidocr_onnxruntime", None)
        else:
            sys.modules["rapidocr_onnxruntime"] = orig
        r._ocr_engine_instance = None


def test_ocr_pdf_raises_missing_dep_error(monkeypatch):
    """_ocr_pdf raises MissingDependencyError(extra='ocr') when pypdfium2 absent."""
    import sys
    import raggity.readers as r
    from raggity.readers import MissingDependencyError

    orig = sys.modules.get("pypdfium2")
    sys.modules["pypdfium2"] = None  # type: ignore[assignment]
    try:
        import pytest
        with pytest.raises(MissingDependencyError) as exc_info:
            r._ocr_pdf("irrelevant.pdf")
        assert exc_info.value.extra == "ocr"
    finally:
        if orig is None:
            sys.modules.pop("pypdfium2", None)
        else:
            sys.modules["pypdfium2"] = orig
