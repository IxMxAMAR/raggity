"""Reader registry for raggity.

Dispatch table: Path.suffix.lower() -> reader function.
Optional-dep readers (docx, html, pptx) do a local import and raise a
friendly RuntimeError if the library is absent.
"""
from __future__ import annotations

import csv
import io
from pathlib import Path

# ---------------------------------------------------------------------------
# Supported extensions
# ---------------------------------------------------------------------------

SUPPORTED_EXTS: set[str] = {
    ".md", ".txt", ".pdf", ".docx", ".html", ".csv", ".pptx",
    # OCR-capable image formats (raggity[ocr] extra)
    ".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp",
}


# ---------------------------------------------------------------------------
# Per-type readers
# ---------------------------------------------------------------------------


def read_txt(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="replace")


def read_md(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="replace")


# ---------------------------------------------------------------------------
# OCR seam — module-level so tests can monkeypatch raggity.readers._run_ocr
# ---------------------------------------------------------------------------

_ocr_engine_instance = None  # lazy singleton; exposed for test teardown


def _run_ocr(src: str) -> str:
    """Run RapidOCR on *src* (image path) and return joined text.

    Raises RuntimeError with install hint if rapidocr_onnxruntime is absent.
    This function is a named module-level seam: tests monkeypatch it directly.
    """
    global _ocr_engine_instance
    try:
        from rapidocr_onnxruntime import RapidOCR  # noqa: PLC0415
    except (ImportError, TypeError):
        raise RuntimeError(
            "OCR needs: pip install raggity[ocr]  "
            "(rapidocr-onnxruntime + pypdfium2)"
        )
    if _ocr_engine_instance is None:
        _ocr_engine_instance = RapidOCR()
    result, _ = _ocr_engine_instance(src)
    if not result:
        return ""
    return "\n".join(line[1] for line in result if line and len(line) > 1)


# ---------------------------------------------------------------------------
# PDF seams — both monkeypatchable for tests
# ---------------------------------------------------------------------------


def _pdf_text(path: str) -> str:
    """Extract embedded text from a PDF via pypdf (existing behaviour)."""
    from pypdf import PdfReader  # already a core dep  # noqa: PLC0415

    reader = PdfReader(path)
    parts = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    return "\n".join(parts)


def _ocr_pdf(path: str) -> str:
    """Render each PDF page to an image via pypdfium2, then OCR via _run_ocr."""
    try:
        import pypdfium2 as pdfium  # noqa: PLC0415
    except ImportError:
        raise RuntimeError(
            "OCR of scanned PDFs needs: pip install raggity[ocr]  "
            "(rapidocr-onnxruntime + pypdfium2)"
        )
    import tempfile, os  # noqa: E401

    doc = pdfium.PdfDocument(path)
    parts: list[str] = []
    with tempfile.TemporaryDirectory() as tmpdir:
        for i, page in enumerate(doc):
            bitmap = page.render(scale=2.0)  # 144 dpi → good for OCR
            pil_image = bitmap.to_pil()
            img_path = os.path.join(tmpdir, f"page_{i}.png")
            pil_image.save(img_path)
            parts.append(_run_ocr(img_path))
    return "\n".join(parts)


def read_pdf(path: str) -> str:
    """Return embedded text; fall back to OCR if the page text is empty."""
    t = _pdf_text(path)
    return t if t.strip() else _ocr_pdf(path)


def read_image(path: str) -> str:
    """OCR an image file via RapidOCR (raggity[ocr] extra required)."""
    return _run_ocr(path)


def read_docx(path: str) -> str:
    try:
        import docx  # python-docx  # noqa: PLC0415
    except ImportError:
        raise RuntimeError("reading .docx needs: pip install raggity[docs]")
    doc = docx.Document(path)
    return "\n".join(para.text for para in doc.paragraphs)


def read_html(path: str) -> str:
    try:
        from bs4 import BeautifulSoup  # beautifulsoup4  # noqa: PLC0415
    except ImportError:
        raise RuntimeError("reading .html needs: pip install raggity[docs]")
    raw = Path(path).read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")
    return soup.get_text(" ", strip=True)


def read_csv(path: str) -> str:
    raw = Path(path).read_text(encoding="utf-8", errors="replace")
    reader = csv.DictReader(io.StringIO(raw))
    rows = []
    for row in reader:
        rows.append(", ".join(f"{k}: {v}" for k, v in row.items()))
    return "\n".join(rows)


def read_pptx(path: str) -> str:
    try:
        from pptx import Presentation  # python-pptx  # noqa: PLC0415
    except ImportError:
        raise RuntimeError("reading .pptx needs: pip install raggity[docs]")
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------

_IMAGE_EXTS: frozenset[str] = frozenset({".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".webp"})

_DISPATCH: dict[str, object] = {
    ".txt": read_txt,
    ".md": read_md,
    ".pdf": read_pdf,
    ".docx": read_docx,
    ".html": read_html,
    ".csv": read_csv,
    ".pptx": read_pptx,
}
# Register image extensions
for _ext in _IMAGE_EXTS:
    _DISPATCH[_ext] = read_image


def read_file(path: str) -> str | None:
    """Read *path* and return its text content, or None for unknown extensions."""
    ext = Path(path).suffix.lower()
    fn = _DISPATCH.get(ext)
    if fn is None:
        return None
    return fn(path)  # type: ignore[operator]
